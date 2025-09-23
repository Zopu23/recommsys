import argparse
import io
import os
from pathlib import Path
from typing import Iterable, List, Optional, Set, Tuple

# Third-party imports are intentionally inside functions where helpful to avoid import errors
# when optional dependencies are missing (e.g., OCR-related packages).


DEFAULT_EXTENSIONS: Tuple[str, ...] = (
	".pdf",
	".docx",
	".pptx",
	".html",
	".htm",
	".png",
	".jpg",
	".jpeg",
	".tif",
	".tiff",
	".bmp",
	".gif",
	".webp",
)


def parse_args() -> argparse.Namespace:
	parser = argparse.ArgumentParser(
		description="Extract text from PDFs, DOCX, PPTX, and images into Markdown files."
	)
	parser.add_argument(
		"--input",
		"-i",
		nargs="+",
		help="One or more input files or directories. If omitted, defaults to data/raw if it exists.",
	)
	parser.add_argument(
		"--output",
		"-o",
		type=Path,
		default=Path("data/ingested"),
		help="Output directory for generated Markdown files (default: data/ingested)",
	)
	parser.add_argument(
		"--recursive",
		"-r",
		action="store_true",
		help="Recurse into subdirectories when inputs are directories.",
	)
	parser.add_argument(
		"--extensions",
		"-e",
		nargs="+",
		help=f"File extensions to include (default: {', '.join(DEFAULT_EXTENSIONS)})",
	)
	parser.add_argument(
		"--ocr",
		action="store_true",
		help="Use OCR for images and as a fallback for PDFs/PPTX images.",
	)
	parser.add_argument(
		"--ocr-lang",
		default="eng",
		help="OCR language (default: eng). Install language data in Tesseract if needed.",
	)
	parser.add_argument(
		"--overwrite",
		action="store_true",
		help="Overwrite existing Markdown outputs if they already exist.",
	)
	parser.add_argument(
		"--min-text-chars",
		type=int,
		default=200,
		help="If extracted PDF text has fewer characters than this and --ocr is set, fallback to OCR (default: 200).",
	)
	return parser.parse_args()


def discover_inputs(inputs: Optional[List[str]], extensions: Optional[List[str]], recursive: bool) -> List[Path]:
	"""
	Return a list of files to process given user inputs and filters.
	"""
	if not inputs:
		default_dir = Path("data/raw")
		if default_dir.exists():
			inputs = [str(default_dir)]
		else:
			raise SystemExit("No --input provided and data/raw does not exist.")

	allowed_exts: Set[str] = set((extensions or DEFAULT_EXTENSIONS))
	allowed_exts = {ext.lower() if ext.startswith(".") else f".{ext.lower()}" for ext in allowed_exts}

	files: List[Path] = []
	for item in inputs:
		p = Path(item)
		if p.is_dir():
			if recursive:
				for root, _, filenames in os.walk(p):
					for name in filenames:
						fp = Path(root) / name
						if fp.suffix.lower() in allowed_exts:
							files.append(fp)
			else:
				for child in p.iterdir():
					if child.is_file() and child.suffix.lower() in allowed_exts:
						files.append(child)
		elif p.is_file():
			if p.suffix.lower() in allowed_exts:
				files.append(p)
			else:
				print(f"[skip] Unsupported extension for file: {p}")
		else:
			print(f"[warn] Path does not exist: {p}")

	files = sorted(set(f.resolve() for f in files))
	if not files:
		raise SystemExit("No matching input files found.")
	return files


def ensure_output_path(output_dir: Path, input_file: Path) -> Path:
	"""
	Return the output Markdown file path for a given input file, mirroring only filenames.
	"""
	output_dir.mkdir(parents=True, exist_ok=True)
	return output_dir / (input_file.stem + ".md")


def ensure_tesseract_available() -> Optional[str]:
	"""
	Try to locate tesseract executable on Windows. Return its path if found, else None.
	"""
	try:
		import shutil
		cmd = shutil.which("tesseract")
		if cmd:
			return cmd
		# Common Windows install path
		candidate = Path("C:/Program Files/Tesseract-OCR/tesseract.exe")
		if candidate.exists():
			return str(candidate)
		return None
	except Exception:
		return None


def extract_text_from_pdf(pdf_path: Path, use_ocr: bool, ocr_lang: str, min_text_chars: int) -> str:
	text = ""
	# First attempt: pdfminer.six text extraction
	try:
		from pdfminer.high_level import extract_text as pdfminer_extract_text
		text = pdfminer_extract_text(str(pdf_path)) or ""
	except Exception as exc:
		print(f"[pdfminer] Failed to extract text from {pdf_path.name}: {exc}")

	if (not text or len(text) < min_text_chars) and use_ocr:
		print(f"[ocr] Falling back to OCR for {pdf_path.name}...")
		try:
			import fitz  # PyMuPDF
			from PIL import Image
			import pytesseract

			cmd = ensure_tesseract_available()
			if cmd:
				pytesseract.pytesseract.tesseract_cmd = cmd

			ocr_pages: List[str] = []
			doc = fitz.open(str(pdf_path))
			for page_index in range(len(doc)):
				page = doc.load_page(page_index)
				pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
				mode = "RGBA" if pix.alpha else "RGB"
				image = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
				page_text = pytesseract.image_to_string(image, lang=ocr_lang)
				ocr_pages.append(f"## Page {page_index + 1}\n\n{page_text.strip()}\n")
			text = f"# {pdf_path.name}\n\n" + "\n".join(ocr_pages)
		except Exception as exc:
			print(f"[ocr] OCR failed for {pdf_path.name}: {exc}")
			# Keep whatever text we may have, even if small

	if text and not text.startswith("# "):
		# Wrap in a basic markdown heading if not already page-scoped
		text = f"# {pdf_path.name}\n\n" + text
	return text or f"# {pdf_path.name}\n\n(Empty or unsupported content)\n"


def extract_text_from_docx(docx_path: Path) -> str:
	try:
		from docx import Document
		doc = Document(str(docx_path))
		parts: List[str] = [f"# {docx_path.name}"]
		for paragraph in doc.paragraphs:
			text = paragraph.text.strip()
			if text:
				parts.append(text)
		# Tables (simple text representation)
		for table in doc.tables:
			parts.append("")
			for row in table.rows:
				cells = [c.text.strip().replace("\n", " ") for c in row.cells]
				if any(cells):
					parts.append(" | ".join(cells))
			parts.append("")
		return "\n\n".join(parts) + "\n"
	except Exception as exc:
		print(f"[docx] Failed to extract from {docx_path.name}: {exc}")
		return f"# {docx_path.name}\n\n(Unable to extract DOCX content)\n"


def extract_text_from_html(html_path: Path) -> str:
	"""
	Extract readable text from HTML using readability (if available) with bs4 fallback.
	"""
	try:
		text_content = ""
		raw_html = html_path.read_text(encoding="utf-8", errors="ignore")
		try:
			from readability import Document  # readability-lxml
			summary_html = Document(raw_html).summary(html_partial=True)
			from bs4 import BeautifulSoup
			soup = BeautifulSoup(summary_html, "lxml")
			# Prefer article-like content
			main = soup
			# Collect paragraphs and headings
			parts: List[str] = []
			title_tag = soup.find("title")
			title_text = (title_tag.get_text(strip=True) if title_tag else html_path.stem)
			parts.append(f"# {title_text}")
			for header in main.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
				h_text = header.get_text(" ", strip=True)
				if h_text:
					level = int(header.name[1]) if header.name and header.name[1:].isdigit() else 2
					parts.append("#" * level + f" {h_text}")
			for p in main.find_all(["p", "li"]):
				p_text = p.get_text(" ", strip=True)
				if p_text:
					parts.append(p_text)
			text_content = "\n\n".join(parts) + "\n"
		except Exception:
			# Fallback: strip all text
			from bs4 import BeautifulSoup
			soup = BeautifulSoup(raw_html, "lxml")
			title_tag = soup.find("title")
			title_text = (title_tag.get_text(strip=True) if title_tag else html_path.stem)
			body_text = soup.get_text("\n", strip=True)
			text_content = f"# {title_text}\n\n{body_text}\n"
		return text_content
	except Exception as exc:
		print(f"[html] Failed to extract from {html_path.name}: {exc}")
		return f"# {html_path.name}\n\n(Unable to extract HTML content)\n"


def _iter_pptx_text_shapes(slide) -> Iterable[str]:
	for shape in slide.shapes:
		try:
			if hasattr(shape, "text_frame") and shape.text_frame is not None:
				yield shape.text_frame.text
		except Exception:
			continue


def _iter_pptx_images(slide) -> Iterable[bytes]:
	try:
		from pptx.enum.shapes import MSO_SHAPE_TYPE
	except Exception:
		return []
	images: List[bytes] = []
	for shape in slide.shapes:
		try:
			if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
				images.append(shape.image.blob)
		except Exception:
			continue
	return images


def extract_text_from_pptx(pptx_path: Path, use_ocr: bool, ocr_lang: str) -> str:
	try:
		from pptx import Presentation
		from PIL import Image
		import pytesseract
		cmd = ensure_tesseract_available()
		if cmd:
			pytesseract.pytesseract.tesseract_cmd = cmd

		prs = Presentation(str(pptx_path))
		parts: List[str] = [f"# {pptx_path.name}"]
		for index, slide in enumerate(prs.slides, start=1):
			parts.append(f"\n## Slide {index}\n")
			for text in _iter_pptx_text_shapes(slide):
				text = (text or "").strip()
				if text:
					parts.append(text)
			if use_ocr:
				for blob in _iter_pptx_images(slide):
					try:
						image = Image.open(io.BytesIO(blob))
						ocr_text = pytesseract.image_to_string(image, lang=ocr_lang).strip()
						if ocr_text:
							parts.append(ocr_text)
					except Exception:
						continue
		return "\n\n".join(parts) + "\n"
	except Exception as exc:
		print(f"[pptx] Failed to extract from {pptx_path.name}: {exc}")
		return f"# {pptx_path.name}\n\n(Unable to extract PPTX content)\n"


def extract_text_from_image(image_path: Path, ocr_lang: str) -> str:
	try:
		from PIL import Image
		import pytesseract
		cmd = ensure_tesseract_available()
		if cmd:
			pytesseract.pytesseract.tesseract_cmd = cmd
		image = Image.open(str(image_path))
		text = pytesseract.image_to_string(image, lang=ocr_lang) or ""
		return f"# {image_path.name}\n\n{text.strip()}\n"
	except Exception as exc:
		print(f"[image] OCR failed for {image_path.name}: {exc}")
		return f"# {image_path.name}\n\n(Unable to OCR image)\n"


def extract_to_markdown(input_file: Path, output_file: Path, use_ocr: bool, ocr_lang: str, min_text_chars: int, overwrite: bool) -> None:
	if output_file.exists() and not overwrite:
		print(f"[skip] {output_file} exists (use --overwrite to replace)")
		return

	print(f"[proc] {input_file} -> {output_file}")
	text = ""
	lower_suffix = input_file.suffix.lower()
	try:
		if lower_suffix == ".pdf":
			text = extract_text_from_pdf(input_file, use_ocr=use_ocr, ocr_lang=ocr_lang, min_text_chars=min_text_chars)
		elif lower_suffix == ".docx":
			text = extract_text_from_docx(input_file)
		elif lower_suffix == ".pptx":
			text = extract_text_from_pptx(input_file, use_ocr=use_ocr, ocr_lang=ocr_lang)
		elif lower_suffix in {".html", ".htm"}:
			text = extract_text_from_html(input_file)
		elif lower_suffix in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif", ".webp"}:
			text = extract_text_from_image(input_file, ocr_lang=ocr_lang)
		else:
			print(f"[warn] Unsupported file type: {input_file}")
			return
		output_file.write_text(text, encoding="utf-8")
	except Exception as exc:
		print(f"[error] Failed to process {input_file.name}: {exc}")


def main() -> None:
	args = parse_args()
	files = discover_inputs(args.input, args.extensions, args.recursive)
	args.output.mkdir(parents=True, exist_ok=True)
	for file_path in files:
		out_path = ensure_output_path(args.output, file_path)
		extract_to_markdown(
			input_file=file_path,
			output_file=out_path,
			use_ocr=args.ocr,
			ocr_lang=args.ocr_lang,
			min_text_chars=args.min_text_chars,
			overwrite=args.overwrite,
		)
	print("Done.")


if __name__ == "__main__":
	main()


